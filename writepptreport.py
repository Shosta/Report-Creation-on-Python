# coding: utf-8
"""
Write the eIoT report in a Powerpoint file that has one slide.\n

"""
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
import variables

def __get_preceding_securityprogress_counters():
    '''
    Get the Security Progress objects counters and return a Tuple that contains them.
    The Security Report from where we are looking for the counters is located in the
    'Reports History' folder. 

    Return:
    A Tuple that has the Security Progress counters in the following order, 
    Not started,
    In progresss,
    Done,
    Not evaluated
    '''
    from datetime import datetime
    date = datetime.now()
    last_month = date.replace(date.year, date.month-1, 1)

    preceding_report_file_name = str(last_month.year) + '-' + str(last_month.month) + variables.SECURITY_REPORT_FILE_NAME
    
    import os.path
    from report_creation_utils import xml_utils
    if os.path.isfile(os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)):
        last_month_security_report_path = os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)
        
        are_notstarted_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'NotStartedObjects', 'counter')
        are_inprogresss_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'InProgressObjects', 'counter')
        are_done_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'DoneObjects', 'counter')
        are_notevaluated_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'NotEvaluatedObjects', 'counter')
        
        
        return [are_notstarted_counter,
                are_inprogresss_counter,
                are_done_counter,
                are_notevaluated_counter]

    return [0, 0, 0, 0]


def __form_securityprogress_chart_data(security_report):
    last_month_securityprogress_counters = __get_preceding_securityprogress_counters()

    chart_data = ChartData()
    chart_data.categories = ['July', 'August']
    chart_data.add_series('Done', (last_month_securityprogress_counters[0], security_report.security_phase_progress.done_objects))
    chart_data.add_series('In Progress', (last_month_securityprogress_counters[1], security_report.security_phase_progress.in_progress_objects))
    chart_data.add_series('Not Started', (last_month_securityprogress_counters[2], security_report.security_phase_progress.not_started_objects))
    chart_data.add_series('Not Eval.', (last_month_securityprogress_counters[3], security_report.security_phase_progress.not_evaluated_objects))

    return chart_data


def __get_preceding_deliveryprocess_counters():
    '''
    Get the delivery security process objects counters and return a Tuple that contains them.
    The Security Report from where we are looking for the counters is located in the
    'Reports History' folder. 

    Return:
    A Tuple that has the Delivery Security process counters in the following order, 
    evaluation,
    qualification,
    security_audit,
    testing,
    in_life
    '''
    from datetime import datetime
    date = datetime.now()
    last_month = date.replace(date.year, date.month-1, 1)

    preceding_report_file_name = str(last_month.year) + '-' + str(last_month.month) + variables.SECURITY_REPORT_FILE_NAME
    
    import os.path
    from report_creation_utils import xml_utils
    if os.path.isfile(os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)):
        last_month_security_report_path = os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)
        
        evaluation_objects_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'EvaluationObjects', 'counter')
        qualification_objects_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'QualificationObjects', 'counter')
        security_audit_objects_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'SecurityAuditObjects', 'counter')
        testing_objects_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'TestingObjects', 'counter')
        in_life_objects_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'InLifeObjects', 'counter')

        return [evaluation_objects_counter,
                qualification_objects_counter,
                security_audit_objects_counter,
                testing_objects_counter,
                in_life_objects_counter]

    return [0, 0, 0, 0, 0]

def __form_deliveryprocess_chart_data(security_report): 
    last_month_delivery_counters = __get_preceding_deliveryprocess_counters()

    chart_data = ChartData()
    chart_data.categories = ['July', 'August']
    chart_data.add_series('Eval.', (last_month_delivery_counters[0], security_report.delivery_security_process.evaluation_objects))
    chart_data.add_series('Qualif.', (last_month_delivery_counters[1], security_report.delivery_security_process.qualification_objects))
    chart_data.add_series('Tests', (last_month_delivery_counters[2], security_report.delivery_security_process.testing_objects))
    chart_data.add_series('Audits', (last_month_delivery_counters[3], security_report.delivery_security_process.security_audit_objects))
    chart_data.add_series('In Life', (last_month_delivery_counters[4], security_report.delivery_security_process.in_life_objects))

    return chart_data


def __add_data_label_on_chart_columns(chart):
    # Add data labels to Chart columns
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LABEL_POSITION

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(13)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END



def __add_legend_to_chart(chart, position):
    chart.has_legend = True
    chart.legend.position = position
    chart.legend.include_in_layout = False


def __write_securityprogress_chart(slide, security_report):
    # Define chart data 
    chart_data = __form_securityprogress_chart_data(security_report)

    # add chart to slide --------------------
    x, y, cx, cy = Cm(9), Cm(11.5), Cm(16.5), Cm(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Add data labels to Chart columns
    __add_data_label_on_chart_columns(chart)

    # Add legend to chart
    from pptx.enum.chart import XL_LEGEND_POSITION
    __add_legend_to_chart(chart, XL_LEGEND_POSITION.RIGHT)


def __write_deliveryprocess_chart(slide, security_report):
    # Define chart data 
    chart_data = __form_deliveryprocess_chart_data(security_report)

    # add chart to slide --------------------
    x, y, cx, cy = Cm(9), Cm(2.75), Cm(16.5), Cm(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Add data labels to Chart columns
    __add_data_label_on_chart_columns(chart)

    # Add legend to chart
    from pptx.enum.chart import XL_LEGEND_POSITION
    __add_legend_to_chart(chart, XL_LEGEND_POSITION.RIGHT)


def __draw_riskiest_objects_linechart(slide, security_report):
    '''
    It draws an horizontal line charts to the Dashboard slide.
    The chart describes the top 5 objects that have the highest number of risks.
    '''
    riskiestobjects_list = security_report.risks.riskiestobjects_list
    length = len(riskiestobjects_list)

    chart_data = ChartData()
    highrisks_chartserie = []
    moderaterisks_chartserie = []
    lightrisks_chartserie = []

    
    # We have less than 5 objects so we take every objets and display it.
    if (length < 5):
        for iot_object in riskiestobjects_list:
            chart_data.categories.add_category(iot_object.object_name)

            highrisks_chartserie.append(iot_object.high_risks_counter)
            moderaterisks_chartserie.append(iot_object.moderate_risks_counter)
            lightrisks_chartserie.append(iot_object.light_risks_counter)
    else:
        chart_data.categories = [riskiestobjects_list[length-1].object_name,
                                riskiestobjects_list[length-2].object_name,
                                riskiestobjects_list[length-3].object_name,
                                riskiestobjects_list[length-4].object_name,
                                riskiestobjects_list[length-5].object_name]

        highrisks_chartserie = [riskiestobjects_list[length-1].high_risks_counter,
                                riskiestobjects_list[length-2].high_risks_counter,
                                riskiestobjects_list[length-3].high_risks_counter,
                                riskiestobjects_list[length-4].high_risks_counter,
                                riskiestobjects_list[length-5].high_risks_counter]
        moderaterisks_chartserie = [riskiestobjects_list[length-1].moderate_risks_counter,
                                    riskiestobjects_list[length-2].moderate_risks_counter,
                                    riskiestobjects_list[length-3].moderate_risks_counter,
                                    riskiestobjects_list[length-4].moderate_risks_counter,
                                    riskiestobjects_list[length-5].moderate_risks_counter]
        lightrisks_chartserie = [riskiestobjects_list[length-1].light_risks_counter,
                                 riskiestobjects_list[length-2].light_risks_counter,
                                 riskiestobjects_list[length-3].light_risks_counter,
                                 riskiestobjects_list[length-4].light_risks_counter,
                                 riskiestobjects_list[length-5].light_risks_counter]

    chart_data.add_series('High Risks', highrisks_chartserie)
    chart_data.add_series('Moderate Risks', moderaterisks_chartserie)
    chart_data.add_series('Light Risks', lightrisks_chartserie)
    
     # add chart to slide --------------------
    x, y, cx, cy = Cm(9), Cm(11.5), Cm(16.5), Cm(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
    ).chart

    __add_data_label_on_chart_columns(chart)

    # Add legend to chart
    from pptx.enum.chart import XL_LEGEND_POSITION
    __add_legend_to_chart(chart, XL_LEGEND_POSITION.BOTTOM)
    

def __get_preceding_risks_counters():
    '''
    Get the High risks and Light risks counters and return a Tuple that contains them.
    The Security Report from where we are looking for the counters is located in the
    'Reports History' folder. 

    Return:
    A Tuple that has the Delivery Security process counters in the following order, 
    High risks,
    Moderate risks,
    Light risks
    '''
    from datetime import datetime
    date = datetime.now()
    last_month = date.replace(date.year, date.month-1, 1)

    preceding_report_file_name = str(last_month.year) + '-' + str(last_month.month) + variables.SECURITY_REPORT_FILE_NAME
    
    import os.path
    from report_creation_utils import xml_utils
    if os.path.isfile(os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)):
        last_month_security_report_path = os.path.join(variables.SECURITY_PROJECTS_DIRPATH, variables.HISTORY_FOLDER_NAME, preceding_report_file_name)
        
        high_risks_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'HighRisks', 'counter')
        moderate_risks_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'ModerateRisks', 'counter')
        light_risks_counter = xml_utils.get_attribute_value(last_month_security_report_path, 'LightRisks', 'counter')
        
        return [high_risks_counter,
                moderate_risks_counter,
                light_risks_counter]

    return [0, 0, 0]


def __write_dashboard_slide(presentation, security_report):
    from SecurityReport import SecurityReport, Risks
    # We choose the first slide of the presentation. Dashboard Slide.
    slide = presentation.slides[0]
    last_month_risks = __get_preceding_risks_counters()

    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            
            if shape.has_text_frame:
                shape.text = '?!!'

                # TODO: change this if else statement with a Python switch case : 
                # https://bytebaker.com/2008/11/03/switch-case-statement-in-python/
                
                # High Risks this month
                if shape.placeholder_format.idx == 10:
                    shape.text = str(security_report.risks.high_risks)

                # High Risks last month
                if shape.placeholder_format.idx == 11:
                    shape.text = str(last_month_risks[0])

                # Moderate Risks this month
                if shape.placeholder_format.idx == 18:
                    shape.text = str(security_report.risks.moderate_risks)

                # Moderate Risks last month
                if shape.placeholder_format.idx == 19:
                    shape.text = str(last_month_risks[1])

                # Light Risks this month
                if shape.placeholder_format.idx == 12:
                    shape.text = str(security_report.risks.light_risks)

                # Light Risks last month
                if shape.placeholder_format.idx == 13:
                    shape.text = str(last_month_risks[1])

                # Total B2B Objects
                if shape.placeholder_format.idx == 14:
                    shape.text = str(security_report.b2b_objects_count)

                # Total B2C Objects
                if shape.placeholder_format.idx == 15:
                    shape.text = str(security_report.b2c_objects_count)

    # Write Chart on slide
    __write_deliveryprocess_chart(slide, security_report)
    #__write_securityprogress_chart(slide, security_report)
    __draw_riskiest_objects_linechart(slide, security_report)


def __write_highlights_slide(presentation, security_report):
    # We choose the second slide of the presentation. Highlights Slide.
    slide = presentation.slides[1]

    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            #print('%d, %s' % (phf.idx, phf.type))

            if shape.has_text_frame:
                shape.text = '?!!'

                # Monthly Highlight
                if shape.placeholder_format.idx == 10:
                    shape.text = security_report.highlights.monthly_highlights

                # Risks
                if shape.placeholder_format.idx == 11:
                    shape.text = security_report.highlights.highlighted_risks

                # Identified Problems
                if shape.placeholder_format.idx == 12:
                    shape.text = security_report.highlights.problems_identified


def write_dashboard_file(security_projects_folder, ppt_file_name, security_report):

    import os.path
    path = os.path.join(os.path.dirname(__file__),  variables.TEMPLATES_FOLDER_NAME, 'SecurityReport.pptx')
    
    # Open the default Powerpoint file we are going to use to write the Report.
    prs = Presentation(path)
    __write_dashboard_slide(prs, security_report)
    __write_highlights_slide(prs, security_report)

    from datetime import datetime
    now = datetime.now()
    ppt_abs_path = os.path.join(security_projects_folder,
                                variables.HISTORY_FOLDER_NAME,
                                str(now.year) + "-" + str(now.month)+ " - " + ppt_file_name)
    prs.save(ppt_abs_path)

    